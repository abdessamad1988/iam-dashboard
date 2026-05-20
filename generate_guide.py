"""
Génère le guide utilisateur IAM Dashboard au format PowerPoint.
Usage : python generate_guide.py
Sortie : IAM_Dashboard_Guide_Utilisateur.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ───────────────────────────────────────────────────────────────────
C_BLUE_DARK  = RGBColor(0x0D, 0x21, 0x37)   # #0D2137
C_BLUE       = RGBColor(0x19, 0x76, 0xD2)   # #1976D2
C_BLUE_LIGHT = RGBColor(0xE3, 0xF2, 0xFD)   # #E3F2FD
C_GREEN      = RGBColor(0x38, 0x8E, 0x3C)   # #388E3C
C_RED        = RGBColor(0xD3, 0x2F, 0x2F)   # #D32F2F
C_ORANGE     = RGBColor(0xF5, 0x7C, 0x00)   # #F57C00
C_PURPLE     = RGBColor(0x7B, 0x1F, 0xA2)   # #7B1FA2
C_GREY       = RGBColor(0x54, 0x6E, 0x7A)   # #546E7A
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BG         = RGBColor(0xF5, 0xF7, 0xFA)   # #F5F7FA

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


def _rgb(r, g, b):
    return RGBColor(r, g, b)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


# ── Primitives ────────────────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill_rgb, line_rgb=None, line_width=Pt(0)):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def textbox(slide, l, t, w, h, text, font_size=Pt(12), bold=False, color=C_WHITE,
            align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = font_size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_bullet_box(slide, l, t, w, h, items, font_size=Pt(11), color=C_BLUE_DARK,
                   bullet="▸ ", line_spacing=1.15):
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf    = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after  = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}{item}"
        run.font.size  = font_size
        run.font.color.rgb = color
    return txBox


def header_bar(slide, title, subtitle="", bg=C_BLUE):
    rect(slide, 0, 0, W, Inches(1.35), bg)
    textbox(slide, Inches(0.4), Inches(0.12), Inches(12), Inches(0.7),
            title, font_size=Pt(28), bold=True, color=C_WHITE)
    if subtitle:
        textbox(slide, Inches(0.4), Inches(0.82), Inches(12), Inches(0.45),
                subtitle, font_size=Pt(13), color=RGBColor(0xBB, 0xDE, 0xFB), italic=True)


def footer(slide, text="IAM Dashboard — Guide Utilisateur"):
    rect(slide, 0, Inches(7.1), W, Inches(0.4), C_BLUE_DARK)
    textbox(slide, Inches(0.3), Inches(7.13), Inches(12.7), Inches(0.35),
            text, font_size=Pt(9), color=RGBColor(0x90, 0xA4, 0xAE), align=PP_ALIGN.CENTER)


def card(slide, l, t, w, h, title, body_lines, title_color=C_BLUE, border_color=C_BLUE):
    rect(slide, l, t, w, h, C_WHITE, border_color, Pt(1.5))
    rect(slide, l, t, w, Inches(0.38), title_color)
    textbox(slide, l + Inches(0.12), t + Inches(0.04), w - Inches(0.24), Inches(0.32),
            title, font_size=Pt(11), bold=True, color=C_WHITE)
    add_bullet_box(slide, l + Inches(0.12), t + Inches(0.44),
                   w - Inches(0.24), h - Inches(0.55),
                   body_lines, font_size=Pt(10), color=C_BLUE_DARK, bullet="• ")


def kpi_box(slide, l, t, val, label, color):
    rect(slide, l, t, Inches(2.8), Inches(1.1), color)
    textbox(slide, l, t + Inches(0.1), Inches(2.8), Inches(0.55),
            str(val), font_size=Pt(28), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    textbox(slide, l, t + Inches(0.65), Inches(2.8), Inches(0.4),
            label, font_size=Pt(10), color=C_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_cover(prs):
    sld = blank_slide(prs)
    # Fond sombre
    rect(sld, 0, 0, W, H, C_BLUE_DARK)
    # Bande bleue gauche
    rect(sld, 0, 0, Inches(0.5), H, C_BLUE)
    # Titre principal
    textbox(sld, Inches(1), Inches(1.4), Inches(11), Inches(1.4),
            "IAM Dashboard", font_size=Pt(52), bold=True, color=C_WHITE)
    # Sous-titre
    textbox(sld, Inches(1), Inches(2.9), Inches(11), Inches(0.7),
            "Guide Utilisateur — Supervision Identity & Access Management",
            font_size=Pt(20), color=RGBColor(0xBB, 0xDE, 0xFB), italic=True)
    # Ligne décorative
    rect(sld, Inches(1), Inches(3.7), Inches(6), Inches(0.06), C_BLUE)
    # Sources
    textbox(sld, Inches(1), Inches(3.95), Inches(11), Inches(0.5),
            "Sources : RH  ·  Active Directory  ·  LDAP  ·  Logs de provisioning",
            font_size=Pt(14), color=RGBColor(0x78, 0x90, 0x9C))
    # Version + date
    textbox(sld, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
            "Version 1.0  |  2026", font_size=Pt(11),
            color=RGBColor(0x54, 0x6E, 0x7A))
    # Icône
    textbox(sld, Inches(10.5), Inches(1.2), Inches(2), Inches(2),
            "🔐", font_size=Pt(90), color=C_WHITE, align=PP_ALIGN.CENTER)


def slide_sommaire(prs):
    sld = blank_slide(prs)
    rect(sld, 0, 0, W, H, C_BG)
    header_bar(sld, "Sommaire", "Vue d'ensemble des sections du guide")
    footer(sld)

    pages = [
        ("01", "Import des Données",    "Upload CSV mensuel — RH, AD, LDAP, Logs",               C_GREEN),
        ("02", "Vue Globale",           "KPIs, évolution mensuelle, alertes rapides",              C_BLUE),
        ("03", "Mobilité RH",           "Extra-entité, intra-site, intra-fonction",                C_PURPLE),
        ("04", "Provisioning",          "Créations, modifications, désactivations, logs",          C_ORANGE),
        ("05", "Qualité des Données",   "Comptes sans mail, entité manquante, incohérences",       C_GREY),
        ("06", "Sécurité",              "Comptes bloqués, MDP expirés, inactifs 90j",              C_RED),
        ("07", "Analyse LDAP ↔ RH",    "Orphelins, désactivations, manquants, doublons (A–F)",    C_BLUE),
        ("08", "Analyse AD ↔ RH",      "Orphelins, bloqués, incohérences lastLogon/timestamp",    C_BLUE_DARK),
    ]
    col_w = Inches(5.9)
    for i, (num, title, desc, color) in enumerate(pages):
        col = i % 2
        row = i // 2
        l = Inches(0.35) + col * Inches(6.65)
        t = Inches(1.55) + row * Inches(1.35)
        rect(sld, l, t, col_w, Inches(1.15), C_WHITE, color, Pt(1.5))
        rect(sld, l, t, Inches(0.55), Inches(1.15), color)
        textbox(sld, l + Inches(0.05), t + Inches(0.3), Inches(0.5), Inches(0.55),
                num, font_size=Pt(16), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        textbox(sld, l + Inches(0.65), t + Inches(0.1), col_w - Inches(0.75), Inches(0.42),
                title, font_size=Pt(13), bold=True, color=color)
        textbox(sld, l + Inches(0.65), t + Inches(0.54), col_w - Inches(0.75), Inches(0.5),
                desc, font_size=Pt(10), color=C_GREY)


def slide_import(prs):
    sld = blank_slide(prs)
    rect(sld, 0, 0, W, H, C_BG)
    header_bar(sld, "📂  Import des Données Mensuelles",
               "Page 00 — Chargement des fichiers CSV sources", C_GREEN)
    footer(sld)

    textbox(sld, Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.4),
            "Uploadez les 5 fichiers CSV pour alimenter le dashboard. Le cache est vidé automatiquement après chaque import.",
            font_size=Pt(11), color=C_GREY, italic=True)

    files = [
        ("👥", "rh_M.csv",   "RH Mois M",      "Employés actifs du mois courant\nColonnes : matricule, nom, prenom,\nentite_affectation, site, action, mois", C_BLUE),
        ("🗂️", "rh_M1.csv",  "RH Mois M-1",    "Snapshot mois précédent\nPermet : détection sorties\net calcul mobilité M vs M-1", C_BLUE),
        ("🖥️", "ad.csv",     "Active Directory","employeeID, displayName, enabled\nlastLogon, lockedOut,\npasswordExpired, mail, timestamp", C_BLUE_DARK),
        ("🔑", "ldap.csv",   "LDAP",            "uid, cn, mail, enabled\nrole (admin/user), department\nlast_logon, description", C_PURPLE),
        ("📋", "logs.csv",   "Logs",            "date, matricule, action\nstatut (success/error)\nmessage", C_ORANGE),
    ]
    cw = Inches(2.45)
    for i, (ico, fname, title, desc, color) in enumerate(files):
        l = Inches(0.25) + i * Inches(2.57)
        t = Inches(2.0)
        rect(sld, l, t, cw, Inches(4.6), C_WHITE, color, Pt(1.5))
        rect(sld, l, t, cw, Inches(0.55), color)
        textbox(sld, l + Inches(0.08), t + Inches(0.06), cw - Inches(0.16), Inches(0.45),
                f"{ico}  {title}", font_size=Pt(11), bold=True, color=C_WHITE)
        textbox(sld, l + Inches(0.1), t + Inches(0.65), cw - Inches(0.2), Inches(0.35),
                fname, font_size=Pt(10), bold=True, color=color, italic=True)
        textbox(sld, l + Inches(0.1), t + Inches(1.05), cw - Inches(0.2), Inches(3.4),
                desc, font_size=Pt(9.5), color=C_BLUE_DARK)

    textbox(sld, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.35),
            "💡  Conseil : uploadez d'abord rh_M1.csv puis rh_M.csv · Encodage UTF-8 avec BOM · Séparateur , ou ;",
            font_size=Pt(10), color=C_GREY, italic=True)


def slide_vue_globale(prs):
    sld = blank_slide(prs)
    rect(sld, 0, 0, W, H, C_BG)
    header_bar(sld, "🏠  Vue Globale — IAM", "Page 01 — Tableau de bord synthétique")
    footer(sld)

    # KPI boxes
    kpis = [
        ("245", "Total Utilisateurs", C_BLUE),
        ("238", "Comptes Actifs AD",  C_GREEN),
        ("24",  "Nouveaux Comptes",   C_PURPLE),
        ("91%", "Taux Automatisation",C_ORANGE),
    ]
    for i, (val, label, color) in enumerate(kpis):
        kpi_box(sld, Inches(0.3) + i * Inches(3.1), Inches(1.5), val, label, color)

    # Graphiques décrits
    rect(sld, Inches(0.3), Inches(2.85), Inches(8.0), Inches(2.8), C_WHITE, C_BLUE, Pt(1))
    textbox(sld, Inches(0.5), Inches(2.95), Inches(7.6), Inches(0.35),
            "📊  Évolution mensuelle des actions RH (Ajouter / Modifier / Désactiver)",
            font_size=Pt(11), bold=True, color=C_BLUE_DARK)
    textbox(sld, Inches(0.5), Inches(3.35), Inches(7.6), Inches(2.2),
            "Graphique à barres groupées par mois et par type d'action.\n"
            "Permet de visualiser la tendance des mouvements de personnel\n"
            "sur la période sélectionnée.",
            font_size=Pt(10), color=C_GREY)

    rect(sld, Inches(8.55), Inches(2.85), Inches(4.5), Inches(2.8), C_WHITE, C_PURPLE, Pt(1))
    textbox(sld, Inches(8.7), Inches(2.95), Inches(4.2), Inches(0.35),
            "🥧  Répartition par action",
            font_size=Pt(11), bold=True, color=C_BLUE_DARK)
    textbox(sld, Inches(8.7), Inches(3.35), Inches(4.2), Inches(2.2),
            "Graphique donut montrant\nla proportion RAS / Modifier\n/ Ajouter / Désactiver.",
            font_size=Pt(10), color=C_GREY)

    # Alertes
    rect(sld, Inches(0.3), Inches(5.85), Inches(12.75), Inches(1.2), C_WHITE, C_RED, Pt(1))
    textbox(sld, Inches(0.5), Inches(5.93), Inches(12.3), Inches(0.35),
            "⚠️  Alertes rapides — 4 indicateurs en temps réel",
            font_size=Pt(11), bold=True, color=C_RED)
    add_bullet_box(sld, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.65),
                   ["🔴 Comptes bloqués (lockedOut AD)",
                    "🟡 Mots de passe expirés  ·  🟡 Comptes sans mail  ·  🔴 Incohérences RH / AD"],
                   font_size=Pt(10), color=C_BLUE_DARK, bullet="")


def slide_mobilite(prs):
    sld = blank_slide(prs)
    rect(sld, 0, 0, W, H, C_BG)
    header_bar(sld, "🔁  Mobilité RH", "Page 02 — Détection des mouvements inter/intra-entité M vs M-1", C_PURPLE)
    footer(sld)

    textbox(sld, Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.35),
            "Comparaison automatique entre le fichier RH Mois M et Mois M-1 sur la base du matricule.",
            font_size=Pt(11), color=C_GREY, italic=True)

    rules = [
        ("Extra-entité\n+ changement site",   "M_entite ≠ M-1_entite\nET M_site ≠ M-1_site",      C_RED),
        ("Extra-entité\nsans changement site", "M_entite ≠ M-1_entite\nET M_site = M-1_site",       C_ORANGE),
        ("Intra\n+ changement site",          "M_entite = M-1_entite\nET M_site ≠ M-1_site",        C_PURPLE),
        ("Intra\n+ changement fonction",      "M_entite = M-1_entite\nET M_fonc ≠ M-1_fonc",        C_BLUE),
        ("Aucun\nchangement",                 "Toutes les valeurs\n= M-1",                           C_GREY),
    ]
    cw = Inches(2.45)
    for i, (title, rule, color) in enumerate(rules):
        l = Inches(0.25) + i * Inches(2.57)
        t = Inches(2.0)
        rect(sld, l, t, cw, Inches(2.5), C_WHITE, color, Pt(1.5))
        rect(sld, l, t, cw, Inches(0.55), color)
        textbox(sld, l + Inches(0.08), t + Inches(0.06), cw - Inches(0.16), Inches(0.5),
                title, font_size=Pt(10), bold=True, color=C_WHITE)
        textbox(sld, l + Inches(0.1), t + Inches(0.65), cw - Inches(0.2), Inches(1.8),
                rule, font_size=Pt(10), color=C_BLUE_DARK)

    textbox(sld, Inches(0.4), Inches(4.7), Inches(5.5), Inches(0.35),
            "Flags transverses (non exclusifs) :", font_size=Pt(11), bold=True, color=C_BLUE_DARK)
    add_bullet_box(sld, Inches(0.4), Inches(5.1), Inches(5.5), Inches(1.8),
                   ["flag_intra_site : même entité + site différent",
                    "flag_intra_sans_site : même entité + même site",
                    "flag_intra_fonction : même entité + fonction différente"],
                   font_size=Pt(10), color=C_BLUE_DARK)

    rect(sld, Inches(6.2), Inches(4.6), Inches(6.8), Inches(2.4), C_WHITE, C_PURPLE, Pt(1))
    textbox(sld, Inches(6.4), Inches(4.7), Inches(6.4), Inches(0.35),
            "📊  Visualisations disponibles", font_size=Pt(11), bold=True, color=C_PURPLE)
    add_bullet_box(sld, Inches(6.4), Inches(5.1), Inches(6.4), Inches(1.8),
                   ["Barres empilées par entité (Extra / Intra)",
                    "Diagrammes donut par type de mobilité",
                    "Tableau détaillé avec colonnes M et M-1",
                    "KPIs : Total mobilité, Extra, Intra, Flags"],
                   font_size=Pt(10), color=C_BLUE_DARK)


def slide_provisioning(prs):
    sld = blank_slide(prs)
    rect(sld, 0, 0, W, H, C_BG)
    header_bar(sld, "⚙️  Provisioning", "Page 03 — Suivi des opérations de création / modification / désactivation", C_ORANGE)
    footer(sld)

    cards_data = [
        ("➕ Créations",      C_GREEN,  ["Nouveaux comptes (action = ajouter)", "Répartition par entité et par site", "Trend mensuel des créations"]),
        ("✏️ Modifications",  C_BLUE,   ["Comptes modifiés (action = modifier)", "Détail des changements par attribut", "Historique sur la période"]),
        ("🚫 Désactivations", C_RED,    ["Employés quittant l'organisation", "Comptes à désactiver dans AD/LDAP", "Alertes sur comptes encore actifs"]),
        ("📋 Logs",           C_ORANGE, ["Journal complet des opérations", "Taux de succès / erreur", "Filtrage par statut et par action"]),
    ]
    cw = Inches(3.0)
    for i, (title, color, items) in enumerate(cards_data):
        l = Inches(0.25) + i * Inches(3.22)
        t = Inches(1.6)
        card(sld, l, t, cw, Inches(5.0), title, items, title_color=color, border_color=color)


def slide_qualite(prs):
    sld = blank_slide(prs)
    rect(sld, 0, 0, W, H, C_BG)
    header_bar(sld, "📋  Qualité des Données", "Page 04 — Détection des anomalies dans les données RH / AD / LDAP", C_GREY)
    footer(sld)

    checks = [
        ("📧", "Comptes sans mail",    "Comptes AD ou LDAP\nsans adresse email renseignée.",       C_ORANGE),
        ("🏢", "Entité manquante",      "Employés RH sans entité\nd'affectation définie.",           C_RED),
        ("🔀", "Incohérences RH / AD", "Département AD ≠ entité RH\npour un même matricule.",       C_RED),
        ("📊", "Indicateurs qualité",  "Score global de qualité\npar source de données.",            C_BLUE),
    ]
    cw = Inches(3.0)
    for i, (ico, title, desc, color) in enumerate(checks):
        l = Inches(0.25) + i * Inches(3.22)
        t = Inches(1.6)
        rect(sld, l, t, cw, Inches(3.0), C_WHITE, color, Pt(1.5))
        rect(sld, l, t, cw, Inches(0.5), color)
        textbox(sld, l + Inches(0.1), t + Inches(0.06), cw - Inches(0.2), Inches(0.4),
                f"{ico}  {title}", font_size=Pt(11), bold=True, color=C_WHITE)
        textbox(sld, l + Inches(0.1), t + Inches(0.6), cw - Inches(0.2), Inches(2.3),
                desc, font_size=Pt(10), color=C_BLUE_DARK)

    rect(sld, Inches(0.3), Inches(4.85), Inches(12.75), Inches(1.9), C_WHITE, C_BLUE, Pt(1))
    textbox(sld, Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.35),
            "📊  Graphiques disponibles", font_size=Pt(11), bold=True, color=C_BLUE)
    add_bullet_box(sld, Inches(0.5), Inches(5.35), Inches(12.3), Inches(1.3),
                   ["Distribution des anomalies par entité",
                    "Évolution mensuelle du taux de qualité",
                    "Top 10 des comptes avec le plus d'anomalies",
                    "Tableau détaillé exportable en CSV"],
                   font_size=Pt(10), color=C_BLUE_DARK)


def slide_securite(prs):
    sld = blank_slide(prs)
    rect(sld, 0, 0, W, H, C_BG)
    header_bar(sld, "🔐  Sécurité", "Page 05 — Surveillance des comptes à risque dans l'AD", C_RED)
    footer(sld)

    risks = [
        ("🔒", "Comptes bloqués",        "lockedOut = True dans l'AD\nIntervention immédiate requise",                  C_RED),
        ("🔑", "MDP expirés",            "passwordExpired = True\nCompte actif avec MDP expiré",                        C_RED),
        ("⏰", "Inactifs > 90 jours",    "lastLogon antérieur à 90j\nCompte actif sans connexion récente",               C_ORANGE),
        ("👤", "Comptes désactivés",     "enabled = False dans l'AD\nPossible désactivation manuelle",                   C_GREY),
    ]
    cw = Inches(3.0)
    for i, (ico, title, desc, color) in enumerate(risks):
        l = Inches(0.25) + i * Inches(3.22)
        t = Inches(1.6)
        rect(sld, l, t, cw, Inches(2.5), C_WHITE, color, Pt(1.5))
        rect(sld, l, t, cw, Inches(0.5), color)
        textbox(sld, l + Inches(0.1), t + Inches(0.06), cw - Inches(0.2), Inches(0.4),
                f"{ico}  {title}", font_size=Pt(11), bold=True, color=C_WHITE)
        textbox(sld, l + Inches(0.1), t + Inches(0.6), cw - Inches(0.2), Inches(1.8),
                desc, font_size=Pt(10), color=C_BLUE_DARK)

    rect(sld, Inches(0.3), Inches(4.3), Inches(12.75), Inches(2.4), C_WHITE, C_RED, Pt(1))
    textbox(sld, Inches(0.5), Inches(4.4), Inches(12.3), Inches(0.35),
            "📊  Visualisations et actions recommandées", font_size=Pt(11), bold=True, color=C_RED)
    add_bullet_box(sld, Inches(0.5), Inches(4.8), Inches(6.0), Inches(1.8),
                   ["Graphique en barres : comptes à risque par entité",
                    "Jauge : score de sécurité global",
                    "Tableau trié par criticité avec export CSV"],
                   font_size=Pt(10), color=C_BLUE_DARK)
    add_bullet_box(sld, Inches(6.8), Inches(4.8), Inches(6.0), Inches(1.8),
                   ["Débloquer les comptes verrouillés (lockedOut)",
                    "Forcer la réinitialisation des MDP expirés",
                    "Désactiver les comptes inactifs > 90 jours"],
                   font_size=Pt(10), color=C_RED, bullet="→ ")


def _analyse_sections_slide(prs, source, page_num, color):
    sld = blank_slide(prs)
    rect(sld, 0, 0, W, H, C_BG)
    header_bar(sld, f"🔍  Analyse {source} — Sections A à F",
               f"Page {page_num} — Sous-menu de navigation intégré dans le sidebar", color)
    footer(sld)

    sections_ldap = [
        ("A", "👻", "Orphelins",    f"Présents dans {source}\nmais absents de RH",               C_RED),
        ("B", "🚫", "À désactiver", "B1 : inactifs > 180 j\nB2 : employés sortis actifs",         C_RED),
        ("C", "❓", "Manquants",    f"Présents dans RH\nmais absents de {source}",                C_ORANGE),
        ("D", "⚠️", "Risque",       "Admins / bloqués\nMDP expirés / jamais utilisés",             C_ORANGE),
        ("E", "📋", "Doublons",     "CN dupliqués (LDAP)\nou mail dupliqué",                       C_BLUE),
        ("F", "⚡", "Incohérences", "Nom · Statut · Département\n· Titre (AD uniquement)",         C_BLUE),
    ]
    cw = Inches(2.08)
    for i, (letter, ico, title, desc, color2) in enumerate(sections_ldap):
        col = i % 3
        row = i // 3
        l = Inches(0.25) + col * Inches(2.2) + (Inches(6.7) if col >= 3 else 0)
        if i >= 3:
            col2 = i - 3
            l = Inches(6.9) + col2 * Inches(2.2)
        else:
            l = Inches(0.25) + i * Inches(2.2)
        t = Inches(1.65)
        rect(sld, l, t, cw, Inches(5.0), C_WHITE, color2, Pt(1.5))
        rect(sld, l, t, cw, Inches(0.55), color2)
        textbox(sld, l + Inches(0.06), t + Inches(0.06), cw - Inches(0.12), Inches(0.45),
                f"{letter}. {ico} {title}", font_size=Pt(11), bold=True, color=C_WHITE)
        textbox(sld, l + Inches(0.1), t + Inches(0.65), cw - Inches(0.2), Inches(1.2),
                desc, font_size=Pt(10), color=C_BLUE_DARK)
        # Exports
        textbox(sld, l + Inches(0.1), t + Inches(1.95), cw - Inches(0.2), Inches(0.3),
                "Actions disponibles :", font_size=Pt(9), bold=True, color=color2)
        add_bullet_box(sld, l + Inches(0.1), t + Inches(2.3), cw - Inches(0.2), Inches(2.5),
                       ["Tableau filtrable", "Export CSV", "Graphiques détaillés"],
                       font_size=Pt(9), color=C_GREY, bullet="• ")

    rect(sld, Inches(0.25), Inches(6.75), Inches(12.8), Inches(0.3),
         RGBColor(0xE3, 0xF2, 0xFD))
    textbox(sld, Inches(0.5), Inches(6.78), Inches(12.3), Inches(0.28),
            f"Navigation : cliquez sur une section dans le sous-menu du sidebar pour afficher le détail correspondant",
            font_size=Pt(9.5), color=C_BLUE, italic=True)


def slide_analyse_ldap(prs):
    _analyse_sections_slide(prs, "LDAP", "06", C_BLUE)


def slide_analyse_ad(prs):
    sld = blank_slide(prs)
    rect(sld, 0, 0, W, H, C_BG)
    header_bar(sld, "🖥️  Analyse AD — Spécificités",
               "Page 07 — lastLogon vs timestamp · 4 sous-sections D · F4 Titre/Fonction", C_BLUE_DARK)
    footer(sld)

    textbox(sld, Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.4),
            "L'analyse AD reprend les mêmes sections A–F que l'analyse LDAP, avec des indicateurs supplémentaires :",
            font_size=Pt(11), color=C_GREY, italic=True)

    specifics = [
        ("📅", "lastLogon vs timestamp",
         "Détecte les comptes créés mais jamais utilisés :\nlastLogon absent OU lastLogon ≈ timestamp (delta < 1 jour)\n→ Section D4 : Jamais utilisés",
         C_BLUE_DARK),
        ("🔒", "Section D enrichie",
         "D1 : Comptes bloqués (lockedOut)\nD2 : MDP expirés (actifs)\nD3 : Inactifs > 90 jours\nD4 : Jamais utilisés depuis création",
         C_RED),
        ("💼", "F4 : Titre / Fonction",
         "Incohérence supplémentaire :\ntitle AD ≠ fonction RH\nPermet de détecter les promotions\nnon répercutées dans l'AD",
         C_ORANGE),
        ("🔑", "Clé de jointure",
         "matricule (RH)\n= employeeID (AD)\nToutes les jointures passent\npar cette correspondance",
         C_PURPLE),
    ]
    cw = Inches(3.0)
    for i, (ico, title, desc, color) in enumerate(specifics):
        l = Inches(0.25) + i * Inches(3.22)
        t = Inches(2.1)
        rect(sld, l, t, cw, Inches(4.5), C_WHITE, color, Pt(1.5))
        rect(sld, l, t, cw, Inches(0.55), color)
        textbox(sld, l + Inches(0.1), t + Inches(0.06), cw - Inches(0.2), Inches(0.45),
                f"{ico}  {title}", font_size=Pt(11), bold=True, color=C_WHITE)
        textbox(sld, l + Inches(0.1), t + Inches(0.65), cw - Inches(0.2), Inches(3.7),
                desc, font_size=Pt(10), color=C_BLUE_DARK)


def slide_workflow(prs):
    sld = blank_slide(prs)
    rect(sld, 0, 0, W, H, C_BG)
    header_bar(sld, "📅  Workflow Mensuel Recommandé",
               "Procédure en 6 étapes à suivre chaque début de mois")
    footer(sld)

    steps = [
        ("1", "Import\nfichiers",      "Uploader rh_M1, rh_M,\nad, ldap, logs\nvia Page 00",           C_GREEN),
        ("2", "Vue Globale",           "Vérifier les KPIs\net les alertes\nrapides (Page 01)",           C_BLUE),
        ("3", "Mobilité",              "Valider les\nmouvements RH\nM vs M-1 (Page 02)",                 C_PURPLE),
        ("4", "Analyse\nLDAP",         "Traiter les sections\nA (orphelins) et B\n(désactiver) en prio", C_ORANGE),
        ("5", "Analyse\nAD",           "Vérifier bloqués,\nMDP expirés,\njamais utilisés",               C_RED),
        ("6", "Export\nrapports",      "Télécharger les CSV\npar section pour\nle rapport mensuel",       C_GREY),
    ]

    arrow_color = RGBColor(0xB0, 0xBE, 0xC5)
    cw = Inches(1.9)
    ch = Inches(4.2)
    for i, (num, title, desc, color) in enumerate(steps):
        l = Inches(0.25) + i * Inches(2.15)
        t = Inches(1.65)
        # Cercle numéro (simulé avec carré arrondi)
        rect(sld, l + Inches(0.7), t, Inches(0.55), Inches(0.55), color)
        textbox(sld, l + Inches(0.7), t, Inches(0.55), Inches(0.55),
                num, font_size=Pt(18), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Carte
        rect(sld, l, t + Inches(0.65), cw, ch - Inches(0.65), C_WHITE, color, Pt(1.5))
        rect(sld, l, t + Inches(0.65), cw, Inches(0.5), color)
        textbox(sld, l + Inches(0.08), t + Inches(0.7), cw - Inches(0.16), Inches(0.42),
                title, font_size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        textbox(sld, l + Inches(0.1), t + Inches(1.25), cw - Inches(0.2), ch - Inches(1.4),
                desc, font_size=Pt(9.5), color=C_BLUE_DARK, align=PP_ALIGN.CENTER)
        # Flèche
        if i < 5:
            textbox(sld, l + Inches(1.95), t + Inches(1.8), Inches(0.25), Inches(0.4),
                    "▶", font_size=Pt(16), color=arrow_color, align=PP_ALIGN.CENTER)


def slide_glossaire(prs):
    sld = blank_slide(prs)
    rect(sld, 0, 0, W, H, C_BG)
    header_bar(sld, "📖  Glossaire", "Termes et abréviations utilisés dans le dashboard")
    footer(sld)

    col1 = [
        ("matricule",        "Identifiant unique RH de l'employé"),
        ("employeeID",       "Identifiant employé dans l'AD (= matricule)"),
        ("uid",              "Identifiant unique dans l'annuaire LDAP (= matricule)"),
        ("enabled",          "Compte actif (True) ou désactivé (False)"),
        ("lastLogon",        "Date de dernière connexion au compte AD"),
        ("last_logon",       "Date de dernière connexion au compte LDAP"),
        ("timestamp",        "Date de création / dernière modification du compte AD"),
        ("lockedOut",        "Compte verrouillé suite à trop de tentatives de connexion"),
        ("passwordExpired",  "Mot de passe expiré selon la politique AD"),
    ]
    col2 = [
        ("action",           "RAS / ajouter / modifier / desactiver (fichier RH)"),
        ("mobilite_type",    "Type de mobilité calculé : extra_avec_site, intra_avec_fonction..."),
        ("orphelin",         "Compte présent dans AD/LDAP mais absent du fichier RH"),
        ("cn",               "Common Name LDAP = Prénom NOM de l'employé"),
        ("displayName",      "Nom affiché dans l'AD = Prénom NOM de l'employé"),
        ("department",       "Entité organisationnelle dans AD/LDAP"),
        ("entite_affectation","Entité d'affectation de l'employé dans le fichier RH"),
        ("role",             "Rôle LDAP : user (standard) ou admin (privilégié)"),
        ("cutoff_180",       "Seuil d'inactivité : 180 jours avant la date courante"),
    ]

    textbox(sld, Inches(0.4), Inches(1.5), Inches(6.2), Inches(0.3),
            "Terme", font_size=Pt(10), bold=True, color=C_BLUE)
    textbox(sld, Inches(2.8), Inches(1.5), Inches(4.0), Inches(0.3),
            "Définition", font_size=Pt(10), bold=True, color=C_BLUE)
    textbox(sld, Inches(7.0), Inches(1.5), Inches(6.2), Inches(0.3),
            "Terme", font_size=Pt(10), bold=True, color=C_BLUE)
    textbox(sld, Inches(9.4), Inches(1.5), Inches(4.0), Inches(0.3),
            "Définition", font_size=Pt(10), bold=True, color=C_BLUE)

    rect(sld, Inches(0.3), Inches(1.82), Inches(6.3), Inches(0.03), C_BLUE)
    rect(sld, Inches(6.95), Inches(1.82), Inches(6.3), Inches(0.03), C_BLUE)

    for i, (term, defn) in enumerate(col1):
        t = Inches(1.95) + i * Inches(0.54)
        bg = C_BLUE_LIGHT if i % 2 == 0 else C_WHITE
        rect(sld, Inches(0.3), t, Inches(6.3), Inches(0.52), bg)
        textbox(sld, Inches(0.4), t + Inches(0.08), Inches(2.3), Inches(0.38),
                term, font_size=Pt(9.5), bold=True, color=C_BLUE)
        textbox(sld, Inches(2.8), t + Inches(0.08), Inches(3.7), Inches(0.38),
                defn, font_size=Pt(9.5), color=C_BLUE_DARK)

    for i, (term, defn) in enumerate(col2):
        t = Inches(1.95) + i * Inches(0.54)
        bg = C_BLUE_LIGHT if i % 2 == 0 else C_WHITE
        rect(sld, Inches(6.95), t, Inches(6.3), Inches(0.52), bg)
        textbox(sld, Inches(7.05), t + Inches(0.08), Inches(2.3), Inches(0.38),
                term, font_size=Pt(9.5), bold=True, color=C_BLUE)
        textbox(sld, Inches(9.45), t + Inches(0.08), Inches(3.7), Inches(0.38),
                defn, font_size=Pt(9.5), color=C_BLUE_DARK)


def slide_fin(prs):
    sld = blank_slide(prs)
    rect(sld, 0, 0, W, H, C_BLUE_DARK)
    rect(sld, 0, 0, Inches(0.5), H, C_BLUE)
    textbox(sld, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
            "Merci d'utiliser le\nIAM Dashboard",
            font_size=Pt(40), bold=True, color=C_WHITE)
    textbox(sld, Inches(1), Inches(3.8), Inches(11), Inches(0.55),
            "Pour toute question, contactez l'équipe IAM.",
            font_size=Pt(16), color=RGBColor(0xBB, 0xDE, 0xFB), italic=True)
    rect(sld, Inches(1), Inches(4.5), Inches(4), Inches(0.06), C_BLUE)
    textbox(sld, Inches(10), Inches(1.8), Inches(2.5), Inches(2),
            "🔐", font_size=Pt(90), color=C_WHITE, align=PP_ALIGN.CENTER)
    textbox(sld, Inches(1), Inches(6.0), Inches(11), Inches(0.4),
            "IAM Dashboard V1  |  RH · Active Directory · LDAP",
            font_size=Pt(11), color=RGBColor(0x54, 0x6E, 0x7A))


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def generate():
    prs = new_prs()

    print("  Slide 01 — Couverture...")
    slide_cover(prs)

    print("  Slide 02 — Sommaire...")
    slide_sommaire(prs)

    print("  Slide 03 — Import données...")
    slide_import(prs)

    print("  Slide 04 — Vue Globale...")
    slide_vue_globale(prs)

    print("  Slide 05 — Mobilite...")
    slide_mobilite(prs)

    print("  Slide 06 — Provisioning...")
    slide_provisioning(prs)

    print("  Slide 07 — Qualite...")
    slide_qualite(prs)

    print("  Slide 08 — Securite...")
    slide_securite(prs)

    print("  Slide 09 — Analyse LDAP sections A-F...")
    slide_analyse_ldap(prs)

    print("  Slide 10 — Analyse AD specificites...")
    slide_analyse_ad(prs)

    print("  Slide 11 — Workflow mensuel...")
    slide_workflow(prs)

    print("  Slide 12 — Glossaire...")
    slide_glossaire(prs)

    print("  Slide 13 — Fin...")
    slide_fin(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "IAM_Dashboard_Guide_Utilisateur.pptx")
    prs.save(out)
    print(f"\n  Fichier genere : {out}")
    return out


if __name__ == '__main__':
    print("IAM Dashboard - Generation du guide utilisateur PowerPoint ...")
    generate()
